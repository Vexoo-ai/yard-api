from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceBgeEmbeddings
from langchain_community.vectorstores import FAISS
import os
from pathlib import Path
import concurrent.futures
from typing import List, Tuple, Optional
import torch
import fitz  # PyMuPDF
import docx
import html2text
import pandas as pd
from bs4 import BeautifulSoup
import re
from pptx import Presentation
from mistralai import Mistral
from dotenv import load_dotenv
import logging

# Load environment variables for API keys
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO,
                    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger("DocumentProcessor")
ALLOWED_FILE_TYPES = [
    'xlsx', 'xls', 'pdf', 'ppt', 'pptx', 'docx', 'doc',
    'txt', 'csv', 'html', 'htm', 'png', 'jpg', 'jpeg',
    'tiff', 'bmp', 'gif'
]
FILE_UPLOAD_LIMIT = int(os.getenv('FILE_UPLOAD_LIMIT', 50))


class DocumentConverter:
    """
    Converts various document formats to Markdown for efficient processing.
    Supports PDF, Word documents, PowerPoint presentations, HTML, text files,
    and scanned documents via OCR (Optical Character Recognition).
    """

    def __init__(self, preserve_tables: bool = True, preserve_images: bool = False, use_ocr: bool = True):
        self.preserve_tables = preserve_tables
        self.preserve_images = preserve_images
        self.use_ocr = use_ocr
        self.html_converter = html2text.HTML2Text()
        self.html_converter.ignore_links = False
        self.html_converter.ignore_images = not preserve_images
        self.html_converter.body_width = 0
        self.html_converter.protect_links = True
        self.html_converter.unicode_snob = True
        self.html_converter.tables = preserve_tables

        # Initialize Mistral client for OCR if needed
        self.mistral_client = None
        if self.use_ocr:
            mistral_api_key = os.getenv("MISTRAL_API_KEY")
            if mistral_api_key:
                self.mistral_client = Mistral(api_key=mistral_api_key)
                logger.info("Mistral client initialized successfully for OCR")
            else:
                logger.warning(
                    "MISTRAL_API_KEY not found. OCR functionality will be disabled.")
                self.use_ocr = False

    def convert(self, file_path: str) -> str:
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        file_ext = os.path.splitext(file_path)[1].lower()

        if file_ext == '.pdf':
            markdown_content = self._convert_pdf(file_path)
            if self.use_ocr and self._needs_ocr(markdown_content, file_path):
                logger.info(
                    f"Document appears to be scanned, using OCR: {file_path}")
                return self._convert_with_ocr(file_path)
            return markdown_content
        elif file_ext in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif']:
            if self.use_ocr:
                logger.info(f"Processing image document with OCR: {file_path}")
                return self._convert_with_ocr(file_path)
            else:
                raise ValueError(
                    "OCR is required for image files but is not enabled")
        elif file_ext in ['.docx', '.doc']:
            return self._convert_word(file_path)
        elif file_ext in ['.pptx', '.ppt']:
            return self._convert_powerpoint(file_path)
        elif file_ext in ['.html', '.htm']:
            return self._convert_html(file_path)
        elif file_ext in ['.txt', '.md', '.markdown']:
            return self._convert_text(file_path)
        elif file_ext in ['.csv', '.xlsx', '.xls']:
            return self._convert_tabular(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")

    def _needs_ocr(self, content: str, file_path: str) -> bool:
        doc = fitz.open(file_path)
        page_count = len(doc)
        doc.close()
        chars_per_page = len(content) / max(1, page_count)
        is_mostly_empty = chars_per_page < 100
        has_encoding_issues = len(re.findall(
            r'[^\x00-\x7F]', content)) > len(content) * 0.3
        return is_mostly_empty or has_encoding_issues

    def _convert_with_ocr(self, file_path: str) -> str:
        if not self.mistral_client:
            logger.error("Mistral client not available for OCR")
            return "ERROR: OCR required but Mistral client not available. Please set MISTRAL_API_KEY."

        try:
            logger.info(f"Processing document with OCR: {file_path}")
            with open(file_path, "rb") as file_content:
                uploaded_file = self.mistral_client.files.upload(
                    file={
                        "file_name": os.path.basename(file_path),
                        "content": file_content,
                    },
                    purpose="ocr"
                )
            logger.info(
                f"File uploaded successfully for OCR: {uploaded_file.id}")

            signed_url = self.mistral_client.files.get_signed_url(
                file_id=uploaded_file.id)

            messages = [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "Please extract all text content from this document and format it as markdown. Include ALL text, preserve structure, tables, and formatting. Don't add any of your own commentary."
                        },
                        {
                            "type": "document_url",
                            "document_url": signed_url.url
                        }
                    ]
                }
            ]

            chat_response = self.mistral_client.chat.complete(
                model="mistral-small-latest",
                messages=messages
            )

            content = chat_response.choices[0].message.content
            logger.info(
                f"Document processed successfully with OCR, extracted {len(content)} characters")
            return content

        except Exception as e:
            logger.error(f"Error processing document with OCR: {str(e)}")
            import traceback
            logger.error(traceback.format_exc())
            return f"ERROR: OCR processing failed: {str(e)}"

    def _convert_pdf(self, file_path: str) -> str:
        markdown_content = []
        doc = fitz.open(file_path)

        for i, page in enumerate(doc):
            text = page.get_text("dict")
            blocks = text["blocks"]
            blocks.sort(key=lambda b: (b["bbox"][1], b["bbox"][0]))

            page_content = []
            in_table = False
            table_rows = []

            for block in blocks:
                if block["type"] == 0:
                    if "lines" not in block:
                        continue

                    is_header = False
                    header_level = 0
                    spans = []
                    for line in block["lines"]:
                        if "spans" in line:
                            spans.extend(line["spans"])

                    if spans:
                        font_size = spans[0]["size"]
                        if font_size > 14:
                            is_header = True
                            if font_size > 20:
                                header_level = 1
                            elif font_size > 16:
                                header_level = 2
                            else:
                                header_level = 3

                    block_text = ""
                    for line in block["lines"]:
                        for span in line.get("spans", []):
                            block_text += span["text"]
                        block_text += "\n"

                    block_text = block_text.strip()
                    if not block_text:
                        continue

                    if is_header and block_text and len(block_text) < 200:
                        page_content.append(
                            f"{'#' * header_level} {block_text}")
                    else:
                        if self.preserve_tables and len(block_text.split()) > 1 and "  " in block_text:
                            table_rows.append(block_text)
                            in_table = True
                        else:
                            if in_table:
                                table_md = self._convert_table_rows_to_markdown(
                                    table_rows)
                                page_content.append(table_md)
                                table_rows = []
                                in_table = False
                            page_content.append(block_text + "\n")

                elif block["type"] == 1 and self.preserve_images:
                    page_content.append(f"![Image on page {i+1}]()")

            if in_table and table_rows:
                table_md = self._convert_table_rows_to_markdown(table_rows)
                page_content.append(table_md)

            markdown_content.append("".join(page_content))
            if i < len(doc) - 1:
                markdown_content.append("\n---\n")

        doc.close()
        return "\n".join(markdown_content)

    def _convert_word(self, file_path: str) -> str:
        doc = docx.Document(file_path)
        markdown_content = []

        for para in doc.paragraphs:
            if not para.text.strip():
                continue
            if para.style.name.startswith('Heading'):
                level = int(para.style.name[-1])
                markdown_content.append(f"{'#' * level} {para.text}\n")
            else:
                markdown_content.append(f"{para.text}\n")

        if self.preserve_tables:
            for table in doc.tables:
                table_md = []
                header_row = [cell.text.strip()
                              for cell in table.rows[0].cells]
                table_md.append("| " + " | ".join(header_row) + " |")
                table_md.append(
                    "| " + " | ".join(["---"] * len(header_row)) + " |")
                for row in table.rows[1:]:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_md.append("| " + " | ".join(row_data) + " |")
                markdown_content.append("\n".join(table_md) + "\n")

        return "\n".join(markdown_content)

    def _convert_powerpoint(self, file_path: str) -> str:
        presentation = Presentation(file_path)
        markdown_content = []
        markdown_content.append(f"# {os.path.basename(file_path)}\n")

        for i, slide in enumerate(presentation.slides):
            markdown_content.append(f"## Slide {i+1}\n")
            if slide.shapes.title:
                markdown_content.append(f"### {slide.shapes.title.text}\n")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if not (hasattr(slide.shapes, "title") and shape == slide.shapes.title):
                        text = shape.text.strip()
                        if len(text) < 100 and "\n" not in text:
                            markdown_content.append(f"#### {text}\n")
                        else:
                            for line in text.split("\n"):
                                if line.strip():
                                    if line.lstrip().startswith("•"):
                                        indent_level = line.find("•") // 2
                                        text_content = line.lstrip(
                                            "• \t").strip()
                                        markdown_content.append(
                                            f"{' ' * indent_level}* {text_content}\n")
                                    else:
                                        markdown_content.append(f"{line}\n")

                if hasattr(shape, "has_table") and shape.has_table:
                    if self.preserve_tables:
                        table = shape.table
                        table_md = []
                        header_row = []
                        for cell in table.rows[0].cells:
                            if hasattr(cell, "text_frame") and hasattr(cell.text_frame, "text"):
                                header_row.append(cell.text_frame.text.strip())
                            else:
                                header_row.append("")
                        table_md.append("| " + " | ".join(header_row) + " |")
                        table_md.append(
                            "| " + " | ".join(["---"] * len(header_row)) + " |")
                        for row_idx in range(1, len(table.rows)):
                            row = table.rows[row_idx]
                            row_data = []
                            for cell in row.cells:
                                if hasattr(cell, "text_frame") and hasattr(cell.text_frame, "text"):
                                    row_data.append(
                                        cell.text_frame.text.strip())
                                else:
                                    row_data.append("")
                            table_md.append("| " + " | ".join(row_data) + " |")
                        markdown_content.append("\n".join(table_md) + "\n")

            if self.preserve_images:
                image_count = sum(
                    1 for shape in slide.shapes if hasattr(shape, "image"))
                if image_count > 0:
                    markdown_content.append(
                        f"\n*This slide contains {image_count} image(s)*\n")

            markdown_content.append("\n---\n")

        return "\n".join(markdown_content)

    def _convert_html(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as file:
            html_content = file.read()
        soup = BeautifulSoup(html_content, 'html.parser')
        for script in soup(["script", "style"]):
            script.extract()
        return self.html_converter.handle(str(soup))

    def _convert_text(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()

    def _convert_tabular(self, file_path: str) -> str:
        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext == '.csv':
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path)

        md_lines = []
        filename = os.path.basename(file_path)
        md_lines.append(f"# {os.path.splitext(filename)[0]}\n")
        headers = df.columns.tolist()
        md_lines.append("| " + " | ".join(str(h) for h in headers) + " |")
        md_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
        for _, row in df.iterrows():
            row_values = [str(val).replace("|", "\\|") for val in row.values]
            md_lines.append("| " + " | ".join(row_values) + " |")
        return "\n".join(md_lines)

    def _convert_table_rows_to_markdown(self, text_rows: List[str]) -> str:
        if not text_rows or len(text_rows) < 2:
            return "\n".join(text_rows)

        normalized_rows = [re.sub(r'\s+', ' ', row.strip())
                           for row in text_rows]
        split_rows = [row.split(' ') for row in normalized_rows]
        num_columns = max(len(row) for row in split_rows)

        md_table = []
        cells = split_rows[0]
        if len(cells) < num_columns:
            cells.extend([''] * (num_columns - len(cells)))
        md_table.append("| " + " | ".join(cells) + " |")
        md_table.append("| " + " | ".join(["---"] * num_columns) + " |")

        for cells in split_rows[1:]:
            if len(cells) < num_columns:
                cells.extend([''] * (num_columns - len(cells)))
            md_table.append("| " + " | ".join(cells) + " |")

        return "\n".join(md_table) + "\n"


class Document:
    """A simple document class to mimic langchain's Document structure"""

    def __init__(self, page_content, metadata=None):
        self.page_content = page_content
        self.metadata = metadata or {}


class DocumentProcessor:
    def __init__(self, use_ocr: bool = True):
        self.model_name = "BAAI/bge-small-en-v1.5"
        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        logger.info(f"Using device: {self.device}")

        self.model_kwargs = {
            "device": self.device,
            "trust_remote_code": True
        }
        self.encode_kwargs = {
            "normalize_embeddings": True,
            "batch_size": 32
        }

        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=10000,
            chunk_overlap=1000
        )

        self.embeddings = HuggingFaceBgeEmbeddings(
            model_name=self.model_name,
            model_kwargs=self.model_kwargs,
            encode_kwargs=self.encode_kwargs
        )

        self.vector_store = None
        self.temp_dir = Path("temp")
        self.temp_dir.mkdir(exist_ok=True)

        self.document_converter = DocumentConverter(
            preserve_tables=True,
            preserve_images=False,
            use_ocr=use_ocr
        )

    def save_file_temporarily(self, file_content, filename) -> Path:
        file_path = self.temp_dir / filename
        file_path.write_bytes(file_content)
        return file_path

    def load_document(self, file_info: Tuple[Path, str]):
        file_path, file_type = file_info
        try:
            if file_type.lower() in ['png', 'jpg', 'jpeg', 'tiff', 'bmp', 'gif']:
                logger.info(f"Processing image file with OCR: {file_path}")
            markdown_content = self.document_converter.convert(str(file_path))
            metadata = {"source": str(file_path), "filetype": file_type}
            return [Document(page_content=markdown_content, metadata=metadata)]
        except Exception as e:
            logger.error(f"Error loading {file_path}: {str(e)}")
            return []

    def process_documents(self, files) -> int:
        file_infos = []

        for file in files:
            file_type = file.name.split('.')[-1].lower()
            file_path = self.save_file_temporarily(file.getbuffer(), file.name)
            file_infos.append((file_path, file_type))

        all_docs = []
        ocr_file_infos = []
        non_ocr_file_infos = []

        for file_path, file_type in file_infos:
            if file_type.lower() in ['png', 'jpg', 'jpeg', 'tiff', 'bmp', 'gif'] or \
               (file_type.lower() == 'pdf' and self._might_need_ocr(file_path)):
                ocr_file_infos.append((file_path, file_type))
            else:
                non_ocr_file_infos.append((file_path, file_type))

        if non_ocr_file_infos:
            with concurrent.futures.ThreadPoolExecutor(max_workers=min(len(non_ocr_file_infos), 4)) as executor:
                doc_futures = list(executor.map(
                    self.load_document, non_ocr_file_infos))
                for docs in doc_futures:
                    if docs:
                        chunks = self.text_splitter.split_documents(docs)
                        all_docs.extend(chunks)

        for file_info in ocr_file_infos:
            docs = self.load_document(file_info)
            if docs:
                chunks = self.text_splitter.split_documents(docs)
                all_docs.extend(chunks)

        for file_path, _ in file_infos:
            try:
                file_path.unlink()
            except Exception as e:
                logger.error(
                    f"Error removing temporary file {file_path}: {str(e)}")

        if all_docs:
            if self.vector_store is None:
                self.vector_store = FAISS.from_documents(
                    all_docs,
                    self.embeddings,
                    distance_strategy="METRIC_INNER_PRODUCT"
                )
            else:
                self.vector_store.add_documents(all_docs)
            return len(all_docs)
        return 0

    def _might_need_ocr(self, file_path: Path) -> bool:
        try:
            doc = fitz.open(str(file_path))
            pages_to_check = min(3, len(doc))
            total_text = ""
            for i in range(pages_to_check):
                page = doc[i]
                total_text += page.get_text()
            doc.close()
            chars_per_page = len(total_text) / max(1, pages_to_check)
            return chars_per_page < 100
        except Exception as e:
            logger.error(f"Error checking if PDF needs OCR: {str(e)}")
            return False

    def search_documents(self, query: str, k: int = 4) -> List:
        if self.vector_store is None:
            raise ValueError("No documents have been processed yet")
        return self.vector_store.similarity_search(query, k=k)

    def get_vector_store(self):
        return self.vector_store
