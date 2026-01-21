from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceBgeEmbeddings
from langchain.vectorstores.faiss import FAISS
import os
from pathlib import Path
import concurrent.futures
from typing import List, Tuple
import torch
import fitz  # PyMuPDF
import docx
import html2text
import pandas as pd
from bs4 import BeautifulSoup
import re
from io import BytesIO
from pptx import Presentation

class DocumentConverter:
    """
    Converts various document formats to Markdown for efficient processing.
    Supports PDF, Word documents, PowerPoint presentations, HTML, and text files.
    """
    
    def __init__(self, preserve_tables: bool = True, preserve_images: bool = False):
        """
        Initialize the document converter
        
        Args:
            preserve_tables: Whether to convert tables to markdown tables
            preserve_images: Whether to extract image references (not implemented for all formats)
        """
        self.preserve_tables = preserve_tables
        self.preserve_images = preserve_images
        self.html_converter = html2text.HTML2Text()
        self.html_converter.ignore_links = False
        self.html_converter.ignore_images = not preserve_images
        self.html_converter.body_width = 0  # No wrapping
        self.html_converter.protect_links = True
        self.html_converter.unicode_snob = True
        self.html_converter.tables = preserve_tables
    
    def convert(self, file_path: str) -> str:
        """
        Convert a document file to markdown based on its extension
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Markdown string representation of the document
        
        Raises:
            ValueError: If the file format is not supported
            FileNotFoundError: If the file does not exist
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_ext = os.path.splitext(file_path)[1].lower()
        
        if file_ext == '.pdf':
            return self._convert_pdf(file_path)
        elif file_ext in ['.docx', '.doc']:
            return self._convert_word(file_path)
        elif file_ext in ['.pptx', '.ppt']:  # Added PowerPoint support
            return self._convert_powerpoint(file_path)
        elif file_ext in ['.html', '.htm']:
            return self._convert_html(file_path)
        elif file_ext in ['.txt', '.md', '.markdown']:
            return self._convert_text(file_path)
        elif file_ext in ['.csv', '.xlsx', '.xls']:
            return self._convert_tabular(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")
    
    def _convert_pdf(self, file_path: str) -> str:
        """Convert PDF to markdown"""
        markdown_content = []
        doc = fitz.open(file_path)
        
        for i, page in enumerate(doc):
            text = page.get_text("dict")
            blocks = text["blocks"]
            
            # Sort blocks by vertical position
            blocks.sort(key=lambda b: (b["bbox"][1], b["bbox"][0]))
            
            page_content = []
            in_table = False
            table_rows = []
            
            for block in blocks:
                if block["type"] == 0:  # Text block
                    if "lines" not in block:
                        continue
                        
                    # Check if this appears to be a header (larger font)
                    is_header = False
                    header_level = 0
                    
                    # Get all the spans in all lines of this block
                    spans = []
                    for line in block["lines"]:
                        if "spans" in line:
                            spans.extend(line["spans"])
                    
                    if spans:
                        # If first span has large font, treat as header
                        font_size = spans[0]["size"]
                        if font_size > 14:
                            is_header = True
                            if font_size > 20:
                                header_level = 1
                            elif font_size > 16:
                                header_level = 2
                            else:
                                header_level = 3
                    
                    # Extract text from spans
                    block_text = ""
                    for line in block["lines"]:
                        for span in line.get("spans", []):
                            block_text += span["text"]
                        block_text += "\n"
                    
                    # Clean up text
                    block_text = block_text.strip()
                    if not block_text:
                        continue
                    
                    # Add header markup if detected
                    if is_header and block_text and len(block_text) < 200:  # Headers shouldn't be too long
                        page_content.append(f"{'#' * header_level} {block_text}")
                    else:
                        # Check if this might be part of a table (heuristic - multiple space-separated values)
                        if self.preserve_tables and len(block_text.split()) > 1 and "  " in block_text:
                            # This might be a table row
                            table_rows.append(block_text)
                            in_table = True
                        else:
                            # If we were in a table but this isn't a table row, end the table
                            if in_table:
                                table_md = self._convert_table_rows_to_markdown(table_rows)
                                page_content.append(table_md)
                                table_rows = []
                                in_table = False
                            
                            # Add paragraph
                            page_content.append(block_text + "\n")
                
                elif block["type"] == 1 and self.preserve_images:  # Image block
                    page_content.append(f"![Image on page {i+1}]()")
            
            # If we have pending table rows at the end of the page
            if in_table and table_rows:
                table_md = self._convert_table_rows_to_markdown(table_rows)
                page_content.append(table_md)
            
            markdown_content.append("".join(page_content))
            
            # Add page separator except for the last page
            if i < len(doc) - 1:
                markdown_content.append("\n---\n")
        
        doc.close()
        return "\n".join(markdown_content)
    
    def _convert_word(self, file_path: str) -> str:
        """Convert Word document to markdown"""
        doc = docx.Document(file_path)
        markdown_content = []
        
        for para in doc.paragraphs:
            # Skip empty paragraphs
            if not para.text.strip():
                continue
                
            # Check if paragraph is a heading
            if para.style.name.startswith('Heading'):
                level = int(para.style.name[-1])
                markdown_content.append(f"{'#' * level} {para.text}\n")
            else:
                # Process text formatting
                text = para.text
                
                # Add paragraph with line break
                markdown_content.append(f"{text}\n")
        
        # Process tables if needed
        if self.preserve_tables:
            for table in doc.tables:
                table_md = []
                
                # Header row
                header_row = []
                for cell in table.rows[0].cells:
                    header_row.append(cell.text.strip())
                table_md.append("| " + " | ".join(header_row) + " |")
                
                # Separator row
                table_md.append("| " + " | ".join(["---"] * len(header_row)) + " |")
                
                # Data rows
                for row in table.rows[1:]:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_md.append("| " + " | ".join(row_data) + " |")
                
                markdown_content.append("\n".join(table_md) + "\n")
        
        return "\n".join(markdown_content)
    
    def _convert_powerpoint(self, file_path: str) -> str:
        """Convert PowerPoint presentation to markdown"""
        presentation = Presentation(file_path)
        markdown_content = []
        
        # Add title
        markdown_content.append(f"# {os.path.basename(file_path)}\n")
        
        # Process each slide
        for i, slide in enumerate(presentation.slides):
            # Add slide header
            markdown_content.append(f"## Slide {i+1}\n")
            
            # Process slide title if it exists
            if slide.shapes.title:
                markdown_content.append(f"### {slide.shapes.title.text}\n")
            
            # Process text content from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Check if this is already the title (to avoid duplication)
                    if not (hasattr(slide.shapes, "title") and shape == slide.shapes.title):
                        # Format based on text frame properties if possible
                        text = shape.text.strip()
                        
                        # Simple heuristic: if text is short and all on one line, it might be a subheading
                        if len(text) < 100 and "\n" not in text:
                            markdown_content.append(f"#### {text}\n")
                        else:
                            # Split by lines and process each line
                            lines = text.split("\n")
                            for line in lines:
                                if line.strip():
                                    # Check if line appears to be a bullet point
                                    if line.lstrip().startswith("•"):
                                        # Convert bullet points
                                        indent_level = line.find("•") // 2  # Approximate indentation level
                                        text_content = line.lstrip("• \t").strip()
                                        markdown_content.append(f"{' ' * indent_level}* {text_content}\n")
                                    else:
                                        markdown_content.append(f"{line}\n")
            
                # Note: For tables and images, we'd need more complex processing
                # If shape contains a table
                if hasattr(shape, "has_table") and shape.has_table:
                    if self.preserve_tables:
                        table = shape.table
                        table_md = []
                        
                        # Create header row from first row
                        header_row = []
                        for cell in table.rows[0].cells:
                            if hasattr(cell, "text_frame") and hasattr(cell.text_frame, "text"):
                                header_row.append(cell.text_frame.text.strip())
                            else:
                                header_row.append("")
                                
                        table_md.append("| " + " | ".join(header_row) + " |")
                        table_md.append("| " + " | ".join(["---"] * len(header_row)) + " |")
                        
                        # Data rows
                        for row_idx in range(1, len(table.rows)):
                            row = table.rows[row_idx]
                            row_data = []
                            for cell in row.cells:
                                if hasattr(cell, "text_frame") and hasattr(cell.text_frame, "text"):
                                    row_data.append(cell.text_frame.text.strip())
                                else:
                                    row_data.append("")
                            table_md.append("| " + " | ".join(row_data) + " |")
                        
                        markdown_content.append("\n".join(table_md) + "\n")
            
            # Add placeholder for images on the slide if needed
            if self.preserve_images:
                image_count = sum(1 for shape in slide.shapes if hasattr(shape, "image"))
                if image_count > 0:
                    markdown_content.append(f"\n*This slide contains {image_count} image(s)*\n")
            
            # Add slide separator
            markdown_content.append("\n---\n")
        
        return "\n".join(markdown_content)
    
    def _convert_html(self, file_path: str) -> str:
        """Convert HTML to markdown"""
        with open(file_path, 'r', encoding='utf-8') as file:
            html_content = file.read()
        
        # Use BeautifulSoup to clean up the HTML first
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.extract()
        
        # Convert to markdown
        return self.html_converter.handle(str(soup))
    
    def _convert_text(self, file_path: str) -> str:
        """Convert plain text or already markdown file"""
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    
    def _convert_tabular(self, file_path: str) -> str:
        """Convert CSV or Excel to markdown table"""
        file_ext = os.path.splitext(file_path)[1].lower()
        
        if file_ext == '.csv':
            df = pd.read_csv(file_path)
        else:  # Excel
            df = pd.read_excel(file_path)
        
        # Convert to markdown table
        md_lines = []
        
        # Add title based on filename
        filename = os.path.basename(file_path)
        md_lines.append(f"# {os.path.splitext(filename)[0]}\n")
        
        # Add the table header
        headers = df.columns.tolist()
        md_lines.append("| " + " | ".join(str(h) for h in headers) + " |")
        md_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
        
        # Add the table rows
        for _, row in df.iterrows():
            row_values = [str(val).replace("|", "\\|") for val in row.values]
            md_lines.append("| " + " | ".join(row_values) + " |")
        
        return "\n".join(md_lines)
    
    def _convert_table_rows_to_markdown(self, text_rows: List[str]) -> str:
        """
        Convert detected table rows from PDF to markdown table format
        Uses heuristic to detect columns from spaces
        """
        if not text_rows or len(text_rows) < 2:
            return "\n".join(text_rows)
        
        # Normalize whitespace to better detect columns
        normalized_rows = []
        for row in text_rows:
            # Replace multiple spaces with a single space for processing
            normalized = re.sub(r'\s+', ' ', row.strip())
            normalized_rows.append(normalized)
        
        # Split by spaces and count the number of cells in each row
        split_rows = [row.split(' ') for row in normalized_rows]
        
        # Use the mode number of columns
        num_columns = max(len(row) for row in split_rows)
        
        # Create the markdown table
        md_table = []
        
        # Header row (first row)
        cells = split_rows[0]
        if len(cells) < num_columns:
            cells.extend([''] * (num_columns - len(cells)))
        md_table.append("| " + " | ".join(cells) + " |")
        
        # Separator row
        md_table.append("| " + " | ".join(["---"] * num_columns) + " |")
        
        # Data rows (remaining rows)
        for cells in split_rows[1:]:
            if len(cells) < num_columns:
                cells.extend([''] * (num_columns - len(cells)))
            md_table.append("| " + " | ".join(cells) + " |")
        
        return "\n".join(md_table) + "\n"


class Document:
    """A simple document class to mimic langchain's Document structure"""
    def __init__(self, page_content, metadata=None):
        self.page_content = page_content
        self.metadata = metadata or {}


class DocumentProcessor:
    def __init__(self):
        self.model_name = "BAAI/bge-small-en-v1.5"
        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        print("------- DEVICE --------")
        print(self.device)
        print("------------------------")
        self.model_kwargs = {
            "device": self.device,
            "trust_remote_code": True
        }
        self.encode_kwargs = {
            "normalize_embeddings": True,
            "batch_size": 32  # Increase batch size for faster processing
        }
        
        # Optimize chunk settings
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=10000,  
            chunk_overlap=1000
        )
        
        self.embeddings = HuggingFaceBgeEmbeddings(
            model_name=self.model_name,
            model_kwargs=self.model_kwargs,
            encode_kwargs=self.encode_kwargs
        )
        
        self.vector_store = None
        self.temp_dir = Path("temp")
        self.temp_dir.mkdir(exist_ok=True)
        
        # Initialize our custom document converter
        self.document_converter = DocumentConverter(preserve_tables=True, preserve_images=False)

    def save_file_temporarily(self, file_content, filename) -> Path:
        """Save uploaded file content temporarily and return the path"""
        file_path = self.temp_dir / filename
        file_path.write_bytes(file_content)
        return file_path

    def load_document(self, file_info: Tuple[Path, str]):
        """Load document based on file type using custom converter"""
        file_path, file_type = file_info
        try:
            # Convert the document to markdown using our custom converter
            markdown_content = self.document_converter.convert(str(file_path))
            
            # Create metadata
            metadata = {
                "source": str(file_path),
                "filetype": file_type
            }
            
            # Create a Document object that mimics langchain Document
            return [Document(page_content=markdown_content, metadata=metadata)]
        except Exception as e:
            print(f"Error loading {file_path}: {str(e)}")
            return []

    def process_documents(self, files) -> int:
        """Process multiple documents and update vector store using parallel processing"""
        file_infos = []
        
        # Save all files first
        for file in files:
            file_type = file.name.split('.')[-1].lower()
            if file_type not in ['xlsx', 'xls', 'pdf', 'ppt', 'pptx', 'docx', 'doc', 'txt', 'csv', 'html', 'htm']:
                continue
                
            file_path = self.save_file_temporarily(file.getbuffer(), file.name)
            file_infos.append((file_path, file_type))

        all_docs = []
        
        # Process documents in parallel
        with concurrent.futures.ThreadPoolExecutor(max_workers=min(len(file_infos), 4)) as executor:
            # Load documents in parallel
            doc_futures = list(executor.map(self.load_document, file_infos))
            
            # Process chunks for each document
            for docs in doc_futures:
                if docs:
                    chunks = self.text_splitter.split_documents(docs)
                    all_docs.extend(chunks)

        # Clean up temporary files
        for file_path, _ in file_infos:
            try:
                file_path.unlink()
            except Exception as e:
                print(f"Error removing temporary file {file_path}: {str(e)}")

        if all_docs:
            # Create or update vector store
            if self.vector_store is None:
                self.vector_store = FAISS.from_documents(
                    all_docs, 
                    self.embeddings,
                    # Additional FAISS optimization parameters
                    distance_strategy="METRIC_INNER_PRODUCT"
                )
            else:
                self.vector_store.add_documents(all_docs)
            
            return len(all_docs)
        return 0

    def search_documents(self, query: str, k: int = 4) -> List:
        """Search documents in vector store"""
        if self.vector_store is None:
            raise ValueError("No documents have been processed yet")
        
        return self.vector_store.similarity_search(
            query, 
            k=k,
            search_kwargs={"k": k + 5}  
        )

    def get_vector_store(self):
        """Return the current vector store"""
        return self.vector_store
